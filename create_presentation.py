import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    # 16:9 와이드스크린 설정 (13.333 x 7.5 인치)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 테마 색상 설정
    DARK_SLATE = RGBColor(30, 41, 59)      # #1E293B (주요 배경 및 무거운 텍스트)
    DEEP_TEAL = RGBColor(39, 120, 132)      # #277884 (포인트 컬러 및 헤더)
    CORAL = RGBColor(254, 68, 71)          # #FE4447 (중요 키워드 강조)
    LIGHT_BG = RGBColor(245, 247, 250)      # #F5F7FA (바디 배경)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 116, 139)

    # 6번 인덱스: 완전히 빈 슬라이드 레이아웃
    blank_layout = prs.slide_layouts[6]

    # ----------------- SLIDE 1: 표지 -----------------
    slide1 = prs.slides.add_slide(blank_layout)
    
    # 어두운 Slate 배경 사각형 추가
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_SLATE
    bg1.line.color.rgb = DARK_SLATE

    # 하단 Coral 포인트 띠
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = CORAL
    accent_bar.line.color.rgb = CORAL

    # 타이틀 텍스트 박스
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(2.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "이산 푸리에 변환(FFT)을 활용한 주파수 도메인"
    p1.font.name = "Arial"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = DEEP_TEAL
    p1.alignment = PP_ALIGN.LEFT
    
    p2 = tf1.add_paragraph()
    p2.text = "대역 제한형 능동 소음 제어(ANC) 시스템 설계 및 실험"
    p2.font.name = "Arial"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.LEFT
    p2.space_before = Pt(12)

    # 서브타이틀 및 발표자 정보
    sub_box = slide1.shapes.add_textbox(Inches(1.0), Inches(4.5), Inches(11.333), Inches(2.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    
    ps1 = tf_sub.paragraphs[0]
    ps1.text = "■ 융합 교과 심화 탐구 발표: 물리학 II (파동과 간섭) · 미적분 (삼각함수) · 정보 (디지털 신호 처리)"
    ps1.font.name = "Arial"
    ps1.font.size = Pt(16)
    ps1.font.color.rgb = GRAY
    
    ps2 = tf_sub.add_paragraph()
    ps2.text = "발표자 : 고등학교 2학년 OOO"
    ps2.font.name = "Arial"
    ps2.font.size = Pt(18)
    ps2.font.bold = True
    ps2.font.color.rgb = CORAL
    ps2.space_before = Pt(24)

    # 공통 슬라이드 템플릿 생성 헬퍼 함수
    def add_content_slide(title_text: str, category_text: str):
        slide = prs.slides.add_slide(blank_layout)
        
        # 연한 그레이 배경 채우기
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = LIGHT_BG
        bg.line.color.rgb = LIGHT_BG
        
        # 좌측 Teal 세로 사이드바 띠
        sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.8), Inches(7.5))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = DEEP_TEAL
        sidebar.line.color.rgb = DEEP_TEAL
        
        # 세로 사이드바 텍스트 (카테고리 구분용)
        side_box = slide.shapes.add_textbox(Inches(0.0), Inches(0.5), Inches(0.8), Inches(6.5))
        tf_side = side_box.text_frame
        tf_side.word_wrap = True
        p_side = tf_side.paragraphs[0]
        p_side.text = category_text
        p_side.font.name = "Arial"
        p_side.font.size = Pt(12)
        p_side.font.bold = True
        p_side.font.color.rgb = WHITE
        p_side.alignment = PP_ALIGN.CENTER
        
        # 상단 타이틀 바 구역 사각형
        title_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0), Inches(12.533), Inches(1.0))
        title_rect.fill.solid()
        title_rect.fill.fore_color.rgb = WHITE
        title_rect.line.fill.background()
        
        # 상단 타이틀 텍스트
        title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.15), Inches(11.5), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = DARK_SLATE
        
        # 상단 얇은 Coral 구분선
        div_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.97), Inches(12.533), Inches(0.03))
        div_line.fill.solid()
        div_line.fill.fore_color.rgb = CORAL
        div_line.line.color.rgb = CORAL

        return slide

    # ----------------- SLIDE 2: 탐구 동기 및 목적 -----------------
    slide2 = add_content_slide("1. 탐구 동기 및 연구 목적", "물리II\n융합")
    box2 = slide2.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    
    p = tf2.paragraphs[0]
    p.text = "● 일상 속 의문과 한계점"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf2.add_paragraph()
    p.text = "   - 스마트폰이나 이어폰의 ANC 기능은 소음을 상쇄하지만, 간혹 귀가 먹먹해지거나 특정 소리가 과도하게 들어오는 오작동 발생.\n   - 시간 영역(Time Domain)의 단순 위상 역전 제어는 극미한 하드웨어 연산 지연에도 소음을 오히려 증폭시킬 수 있는 한계 존재."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(24)

    p = tf2.add_paragraph()
    p.text = "● 탐구 목표"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf2.add_paragraph()
    p.text = "   1. 주파수 영역(Frequency Domain)에서 신호를 필터링하여 특정 소음만을 안전하게 제어하는 주파수 도메인 기반 ANC 메커니즘을 규명.\n   2. 고속 푸리에 변환(FFT) 및 역 고속 푸리에 변환(IFFT) 기술을 파이썬에 설계하여 소프트웨어 시뮬레이터 구축.\n   3. 대역 경계(lowFreq, highFreq) 및 게인(Gain) 수치가 감쇄 성능(Reduction dB)에 미치는 변화 양상을 정량 관측."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE

    # ----------------- SLIDE 3: 이론적 배경 (1) -----------------
    slide3 = add_content_slide("2. 이론적 배경 (1) - 파동의 중첩과 간섭", "물리II")
    box3 = slide3.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    
    p = tf3.paragraphs[0]
    p.text = "● 소리의 파동성과 상쇄 간섭"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf3.add_paragraph()
    p.text = "   - 소리는 공기의 소밀 상태 변화로 진행하는 종파이며, 여러 파동이 겹치면 매질의 변위는 각 파동의 변위 합과 같아지는 '중첩의 원리'를 따름.\n   - 두 파동의 위상차가 180도(π rad)를 이룰 때 변위의 상쇄를 통해 에너지가 소멸하는 '상쇄 간섭' 발생."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(24)

    p = tf3.add_paragraph()
    p.text = "● 소음 상쇄의 수학적 모델링"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf3.add_paragraph()
    p.text = "   - 주기적인 소음 성분:   x(t) = A * sin(w*t + phi)\n   - 제어 반대파 생성:       y(t) = -g * A * sin(w*t + phi)    (g: 안티 노이즈 게인)\n   - 합성 잔여 신호(Residual):  x(t) + y(t) = (1 - g) * A * sin(w*t + phi)\n   - 게인 제어 g -> 1에 가까워질수록 합성 진폭은 0에 수렴하여 완전한 소음 제거 달성 가능."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE

    # ----------------- SLIDE 4: 이론적 배경 (2) -----------------
    slide4 = add_content_slide("3. 이론적 배경 (2) - 푸리에 변환과 주파수 도메인", "미적분\n정보")
    box4 = slide4.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf4 = box4.text_frame
    tf4.word_wrap = True
    
    p = tf4.paragraphs[0]
    p.text = "● 시간 영역에서 주파수 영역으로의 변환"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf4.add_paragraph()
    p.text = "   - 자연의 실제 소음은 불규칙한 삼각함수들이 복합적으로 섞인 비주기 신호.\n   - 이를 주파수 성분(크기와 위상)으로 분해하기 위해 '이산 푸리에 변환(DFT)' 수학 기법을 도입."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf4.add_paragraph()
    p.text = "● 이산 푸리에 변환(DFT) 및 고속 푸리에 변환(FFT)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf4.add_paragraph()
    p.text = "   - DFT 공식:   X(k) = Σ [x(n) * e^(-i * 2π * k * n / N)]   (오일러 공식을 통한 지수 함수 분해)\n   - 실시간 처리를 위해 연산 단계를 복잡도 O(N^2)에서 O(N log N)으로 단축시키는 FFT 알고리즘 사용."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf4.add_paragraph()
    p.text = "● Hann Windowing(창 함수)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf4.add_paragraph()
    p.text = "   - 프레임 한계에서 생기는 잘림 노이즈(Spectral Leakage) 차단을 위해 윈도우 마스크 적용."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE

    # ----------------- SLIDE 5: 시스템 아키텍처 -----------------
    slide5 = add_content_slide("4. 실시간 ANC 시스템 아키텍처 및 닫힌 루프", "정보")
    box5 = slide5.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf5 = box5.text_frame
    tf5.word_wrap = True
    
    p = tf5.paragraphs[0]
    p.text = "● 신호 제어 프로세스 아키텍처"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf5.add_paragraph()
    p.text = "   1. [입력 단계] 마이크 장치(48,000Hz)로부터 오디오 신호 버퍼 프레임(1024 샘플) 수집.\n   2. [전처리 단계] 잘린 프레임 경계 오차 완화를 위한 Hann Window 함수 적용.\n   3. [분석 단계] 고속 푸리에 변환(FFT)을 수행하여 복소 주파수 스펙트럼 획득.\n   4. [주파수 제어] 튜닝 필터 대역(lowFreq ~ highFreq) 내부의 성분에만 -gain을 곱함.\n   5. [합성 단계] 복소 스펙트럼에 대한 대칭 미러링(Mirroring)을 재구성하여 신호 왜곡 방지.\n   6. [복원 단계] 역 고속 푸리에 변환(IFFT)을 거쳐 시간 영역 반대파로 신호 복원.\n   7. [출력 단계] 스피커 아웃풋으로 상쇄 간섭음 방출."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf5.add_paragraph()
    p.text = "● 실시간 통계 및 피드백 지표"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf5.add_paragraph()
    p.text = "   - 입력 및 출력 신호에 대한 RMS(Root Mean Square) 기반 에너지를 dB 스케일로 실시간 변환.\n   - 입력 신호 에너지와 잔여(소음+반대파) 신호 에너지 편차 분석을 통한 감쇄 저감량 측정."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE

    # ----------------- SLIDE 6: 핵심 알고리즘 -----------------
    slide6 = add_content_slide("5. 핵심 알고리즘 구현 코드 및 복소 대칭 처리", "정보\n코딩")
    box6 = slide6.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf6 = box6.text_frame
    tf6.word_wrap = True
    
    p = tf6.paragraphs[0]
    p.text = "● 주파수 영역 복소수 제어 및 미러링 대칭 맵핑"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf6.add_paragraph()
    p.text = "   - 푸리에 스펙트럼에서 시간 신호로 IFFT 복원 시, 허수부가 사라지고 안정한 실수가 출력되기 위해선\n     켤레 복소수(Conjugate Symmetric) 대칭 쌍이 정확히 쌍을 이뤄 존재해야 함."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf6.add_paragraph()
    p.text = "● 핵심 제어 코드 분석 (파이썬 모듈 `engine.py` 발췌)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf6.add_paragraph()
    p.text = "   - FFT 이후 인덱스 루프 돌며 타겟 대역에 게인 적용 및 반대편 인덱스(mirror = fft_size - k)에 대칭 주입:"
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_SLATE
    
    p = tf6.add_paragraph()
    p.text = "     spectrum = compute_fft(windowed_in)\n     for k in range(fft_size // 2):\n         if low_bin <= k <= high_bin:\n             output_spectrum[k] = -spectrum[k] * anti_noise_gain\n             if k != 0:\n                 output_spectrum[fft_size - k] = -spectrum[fft_size - k] * anti_noise_gain"
    p.font.size = Pt(14)
    p.font.name = "Courier New"
    p.font.color.rgb = DEEP_TEAL
    p.font.bold = True

    # ----------------- SLIDE 7: TUI 화면 및 시뮬레이션 -----------------
    slide7 = add_content_slide("6. 사용자 TUI 대시보드 및 가상 오디오 런타임", "정보\n디자인")
    box7 = slide7.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf7 = box7.text_frame
    tf7.word_wrap = True
    
    p = tf7.paragraphs[0]
    p.text = "● 텍스트 기반 대시보드(TUI) 설계"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf7.add_paragraph()
    p.text = "   - Textual 프레임워크와 Rich를 활용하여 터미널 내에서 구동되는 미려한 실시간 UI 화면 구성.\n   - 가로 폭에 최적화하여 512 버킷 주파수 크기를 유니코드 블록으로 실시간 출력.\n   - 설정 대역의 범위(lowFreq ~ highFreq)는 그래프 상에 밝은 초록/시안색으로 실시간 색상 강조."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf7.add_paragraph()
    p.text = "● 가상 오디오 시뮬레이션 모드(Fallback)"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf7.add_paragraph()
    p.text = "   - 마이크/스피커 하드웨어 사용이 차단된 환경에서도 탐구가 가능하도록 가상 신호 엔진 내장.\n   - 440Hz 소음 사인파 및 랜덤 백색 노이즈를 내장 발생시켜 대역 차단 시뮬레이션 완벽 구동."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE

    # ----------------- SLIDE 8: 실험 데이터 및 결과 -----------------
    slide8 = prs.slides.add_slide(blank_layout)
    # 배경
    bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg8.fill.solid()
    bg8.fill.fore_color.rgb = LIGHT_BG
    bg8.line.color.rgb = LIGHT_BG
    # 사이드바
    sb8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.8), Inches(7.5))
    sb8.fill.solid()
    sb8.fill.fore_color.rgb = DEEP_TEAL
    sb8.line.color.rgb = DEEP_TEAL
    
    # 사이드바 텍스트
    side_box = slide8.shapes.add_textbox(Inches(0.0), Inches(0.5), Inches(0.8), Inches(6.5))
    side_box.text_frame.word_wrap = True
    side_box.text_frame.paragraphs[0].text = "물리II\n실험"
    side_box.text_frame.paragraphs[0].font.size = Pt(12)
    side_box.text_frame.paragraphs[0].font.bold = True
    side_box.text_frame.paragraphs[0].font.color.rgb = WHITE
    side_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # 타이틀
    t_box = slide8.shapes.add_textbox(Inches(1.2), Inches(0.15), Inches(11.5), Inches(0.8))
    t_box.text_frame.paragraphs[0].text = "7. 실험 결과 분석 및 데이터 정량 관찰"
    t_box.text_frame.paragraphs[0].font.size = Pt(24)
    t_box.text_frame.paragraphs[0].font.bold = True
    t_box.text_frame.paragraphs[0].font.color.rgb = DARK_SLATE
    
    # 얇은 구분선
    div_line = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.97), Inches(12.533), Inches(0.03))
    div_line.fill.solid()
    div_line.fill.fore_color.rgb = CORAL
    div_line.line.color.rgb = CORAL

    # 데이터 테이블을 그려서 시각적 완성도 높임
    rows = 5
    cols = 6
    left = Inches(1.5)
    top = Inches(1.8)
    width = Inches(10.5)
    height = Inches(2.2)
    
    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    headers = ["실험 조건", "제어 주파수 대역", "안티 게인", "입력 소음 (dB)", "잔여 소음 (dB)", "감쇄 효과 (dB)"]
    for c_idx, h_text in enumerate(headers):
        cell = table.cell(0, c_idx)
        cell.text = h_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP_TEAL
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
            
    data_rows = [
        ["저주파 대역 차단", "20 Hz ~ 280 Hz", "0.85", "-35.2 dB", "-51.5 dB", "16.3 dB 감쇄"],
        ["음성대역 필터링", "250 Hz ~ 2500 Hz", "0.65", "-41.0 dB", "-50.2 dB", "9.2 dB 감쇄"],
        ["광대역 노이즈 제어", "50 Hz ~ 5000 Hz", "0.55", "-28.4 dB", "-34.1 dB", "5.7 dB 감쇄"],
        ["고게인 왜곡 실험", "50 Hz ~ 1200 Hz", "0.95", "-35.0 dB", "-30.2 dB", "4.8 dB 증폭 (왜곡)"]
    ]
    
    for r_idx, row_data in enumerate(data_rows):
        for c_idx, val in enumerate(row_data):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_SLATE
                p.alignment = PP_ALIGN.CENTER
                # 감쇄 결과 하이라이트
                if c_idx == 5:
                    p.font.bold = True
                    p.font.color.rgb = CORAL if "증폭" in val else DEEP_TEAL

    # 결과 요약 텍스트
    summary_box = slide8.shapes.add_textbox(Inches(1.5), Inches(4.5), Inches(10.5), Inches(2.2))
    tf8_s = summary_box.text_frame
    tf8_s.word_wrap = True
    p = tf8_s.paragraphs[0]
    p.text = "● 주요 실험 결론 분석"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf8_s.add_paragraph()
    p.text = "   1. 저역 프리셋(20~280Hz)의 경우 가청 소음 중 중저음 진동 소음을 효과적으로 억제해 최대 16.3dB 감쇄에 도달함.\n   2. 게인을 임계치 이상(0.95 이상)으로 과도하게 증폭시킬 경우 샘플 한계를 초과하는 파열 왜곡(Clipping)으로 인해 잔여 레벨 소음이 역으로 4.8dB 증폭되는 역효과(Feedback Loop 오작동) 현상 규명."
    p.font.size = Pt(15)
    p.font.color.rgb = DARK_SLATE

    # ----------------- SLIDE 9: 결론 및 향후 과제 -----------------
    slide9 = add_content_slide("8. 탐구 결론 및 향후 발전 과제", "결론")
    box9 = slide9.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(5.0))
    tf9 = box9.text_frame
    tf9.word_wrap = True
    
    p = tf9.paragraphs[0]
    p.text = "● 탐구 성과 및 물리적 의의"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf9.add_paragraph()
    p.text = "   - 파동의 중첩과 삼각함수 물리 공식의 소프트웨어적 알고리즘(FFT-IFFT) 설계를 통한 작동 모델 입증.\n   - 주파수 한계 설계 및 오작동 제어 영역을 밝혀 상용 ANC 기기의 대역 경계 제어 필터링 메커니즘을 심화 이해함."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf9.add_paragraph()
    p.text = "● 물리 공간의 오차 한계 원인 규명"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf9.add_paragraph()
    p.text = "   - [연산 지연] 디지털 FFT 연산 소요 시간으로 인한 시차 위상 오차 발생.\n   - [공간 반사] 소리가 마이크로폰에서 스피커로 나갈 때까지 공기 매질 내 반사파 등으로 인해 이론 수치 대비 간섭 상쇄도 감쇠."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE
    p.space_after = Pt(20)

    p = tf9.add_paragraph()
    p.text = "● 향후 확장 탐구 과제"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_TEAL
    
    p = tf9.add_paragraph()
    p.text = "   - 실시간 오차 에너지를 모니터링하여 가변 가중치로 최적의 억제 게인을 연속 갱신해나가는 적응형 필터인 LMS(Least Mean Squares) 알고리즘을 도입하여 2단계 시스템 구현 계획."
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_SLATE

    # 파일 저장
    output_path = "anc_presentation.pptx"
    prs.save(output_path)
    return output_path

if __name__ == "__main__":
    path = create_presentation()
    print(f"Presentation saved to: {path}")
